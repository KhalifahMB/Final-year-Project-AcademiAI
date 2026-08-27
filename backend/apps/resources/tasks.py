"""
Document ingestion pipeline (queue: ingestion).

Upload already stored in object storage; this task validates, extracts text,
chunks, embeds (pgvector) and marks the resource ready. All DB work runs
inside tenant_scope() so RLS policies apply.
"""
import logging
import re

from celery import shared_task
from django.db import transaction

logger = logging.getLogger(__name__)


def _simple_chunk(text: str, max_chars: int = 1200, overlap: int = 150) -> list[str]:
    text = re.sub(r"\s+", " ", text or "").strip()
    if not text:
        return []
    chunks = []
    i = 0
    while i < len(text):
        chunks.append(text[i : i + max_chars])
        i += max_chars - overlap
    return chunks


def _extract_text(raw: bytes, content_type: str, filename: str = "") -> str:
    """Extract plain text from supported document types.

    Supported: PDF (pypdf), DOCX (python-docx), PPTX (python-pptx),
    plain text / UTF-8. Returns "" when nothing can be extracted — the
    caller decides whether that is a failure.
    """
    ct = (content_type or "").lower()
    name = (filename or "").lower()

    if raw[:4] == b"%PDF" or "pdf" in ct or name.endswith(".pdf"):
        try:
            from io import BytesIO

            from pypdf import PdfReader

            reader = PdfReader(BytesIO(raw))
            return "\n".join((page.extract_text() or "") for page in reader.pages)
        except Exception as pdf_exc:
            logger.warning("PDF extract failed: %s", pdf_exc)
            return ""

    # OOXML formats share the PK zip header; dispatch by extension/content type.
    is_docx = (
        "wordprocessingml" in ct
        or name.endswith(".docx")
        or raw[:2] == b"PK"
        and name.endswith(".docx")
    )
    if is_docx:
        try:
            from io import BytesIO

            import docx

            document = docx.Document(BytesIO(raw))
            parts = [p.text for p in document.paragraphs]
            for table in document.tables:
                for row in table.rows:
                    parts.append("\t".join(cell.text for cell in row.cells))
            return "\n".join(parts)
        except Exception as docx_exc:
            logger.warning("DOCX extract failed: %s", docx_exc)
            return ""

    is_pptx = (
        "presentationml" in ct
        or name.endswith(".pptx")
        or raw[:2] == b"PK"
        and name.endswith(".pptx")
    )
    if is_pptx:
        try:
            from io import BytesIO

            from pptx import Presentation

            presentation = Presentation(BytesIO(raw))
            parts = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            parts.append(
                                "\t".join(cell.text for cell in row.cells)
                            )
            return "\n".join(parts)
        except Exception as pptx_exc:
            logger.warning("PPTX extract failed: %s", pptx_exc)
            return ""

    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        return ""


@shared_task(bind=True, max_retries=3, default_retry_delay=60)
def process_resource_ingestion(self, resource_id: str, version_id: str, tenant_id: str):
    from django.conf import settings

    from apps.common.ai import generate_embeddings
    from apps.common.db import tenant_scope
    from apps.common.security.file_validation import FileValidationError, validate_upload_bytes
    from apps.common.storage import get_s3_client
    from apps.resources.models import Resource, ResourceChunk, ResourceVersion

    with tenant_scope(tenant_id):
        # Resolve rows under RLS; a missing pair aborts quietly.
        resource = Resource.objects.filter(id=resource_id).first()
        if resource is None:
            return {"status": "failed", "error": "resource not found"}
        version = ResourceVersion.objects.filter(id=version_id, resource=resource).first()
        if version is None:
            return {"status": "failed", "error": "version not found"}

        resource.processing_status = Resource.ProcessingStatus.PROCESSING
        resource.processing_error = ""
        resource.save(update_fields=["processing_status", "processing_error", "updated_at"])

    tenant_id = resource.tenant_id
    try:
        client = get_s3_client()
        obj = client.get_object(Bucket=settings.AWS_STORAGE_BUCKET_NAME, Key=version.storage_key)
        raw = obj["Body"].read()
        content_type = obj.get("ContentType") or ""

        try:
            validate_upload_bytes(raw, content_type=content_type, filename=version.storage_key or "")
        except FileValidationError as fv:
            raise ValueError(str(fv))

        text = _extract_text(raw, content_type, filename=version.storage_key or "")
        parts = []
        has_text = bool(text.strip())
        if has_text:
            parts = _simple_chunk(text)
        else:
            logger.info("No extractable text found for resource=%s, marking as ready with 0 chunks.", resource_id)

        # Embeddings are an external call — do them before opening the write scope.
        embeddings = generate_embeddings(parts) if parts else []

        with tenant_scope(tenant_id), transaction.atomic():
            ResourceChunk.objects.filter(resource_version=version).delete()
            for idx, (content, emb) in enumerate(zip(parts, embeddings)):
                ResourceChunk.objects.create(
                    tenant_id=tenant_id,
                    resource_version=version,
                    chunk_index=idx,
                    content=content,
                    embedding=emb,
                    token_count=len(content.split()),
                    metadata={},
                )
            resource.storage_key = version.storage_key
            resource.processing_status = Resource.ProcessingStatus.READY
            resource.has_extractable_text = has_text
            # Persist what storage reported so preview can classify the file
            # without re-downloading it.
            if content_type and not resource.mime_type:
                resource.mime_type = content_type
                resource.save(
                    update_fields=[
                        "storage_key", "processing_status", "mime_type",
                        "has_extractable_text", "updated_at",
                    ]
                )
            else:
                resource.save(
                    update_fields=["storage_key", "processing_status", "has_extractable_text", "updated_at"]
                )

        logger.info("Ingestion complete resource=%s chunks=%s", resource_id, len(parts))
        return {"status": "completed", "chunks": len(parts)}
    except Exception as exc:
        logger.exception("Ingestion failed resource=%s", resource_id)
        with tenant_scope(tenant_id):
            resource.processing_status = Resource.ProcessingStatus.FAILED
            resource.processing_error = str(exc)[:2000]
            resource.save(update_fields=["processing_status", "processing_error", "updated_at"])
        raise self.retry(exc=exc)
