"""Text extraction tests for the ingestion pipeline."""
import io

import pytest

from apps.resources.tasks import _extract_text, _simple_chunk


def _make_docx(paragraphs):
    import docx

    document = docx.Document()
    for text in paragraphs:
        document.add_paragraph(text)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _make_pptx(slides):
    from pptx import Presentation

    presentation = Presentation()
    for lines in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        frame = slide.placeholders[1].text_frame
        frame.text = lines[0]
        for extra in lines[1:]:
            frame.add_paragraph().text = extra
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_extract_docx_paragraphs():
    raw = _make_docx(["Introduction to Algorithms", "Greedy methods are covered in week 3."])
    text = _extract_text(raw, "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "lecture.docx")
    assert "Introduction to Algorithms" in text
    assert "Greedy methods" in text


def test_extract_pptx_slides():
    raw = _make_pptx([["Big-O Notation", "O(n log n) for merge sort"]])
    text = _extract_text(raw, "application/vnd.openxmlformats-officedocument.presentationml.presentation", "slides.pptx")
    assert "Big-O Notation" in text
    assert "merge sort" in text


def test_extract_plain_text():
    assert _extract_text("hello world".encode(), "text/plain", "notes.txt") == "hello world"


def test_extract_binary_gives_empty():
    assert _extract_text(b"\x00\x01\x02\xff\xfe", "application/octet-stream", "blob.bin") == ""


def test_simple_chunk_overlap():
    text = "word " * 1000  # 5000 chars
    chunks = _simple_chunk(text.strip(), max_chars=1200, overlap=150)
    assert len(chunks) >= 4
    # Consecutive chunks overlap
    assert chunks[0][-50:] == chunks[1][:50]


def test_simple_chunk_empty():
    assert _simple_chunk("") == []
    assert _simple_chunk("   \n\t ") == []
